from pptx import Presentation
from pptx.util import Pt
import pandas as pd
import math

def round_up(number):
    return int(math.ceil(number))

def grunteco(tonns,burt_length,density,weeks,burt_wall, euro):
    burt_width = 8
    burt_wall = int(burt_wall)
    if burt_wall == 1:
        burt_hieght = 3.5
    else:
        burt_hieght = 2.5

    burt_volume = (burt_width*burt_wall+((0.25*burt_width)*(burt_hieght-burt_wall))+(((0.75*burt_width)* (burt_hieght-burt_wall))*0.5))*(burt_length-3)
    total_volume = tonns / density
    total_burts = round_up(total_volume/(52/weeks)/burt_volume)
    square = total_burts * (burt_length+1) * (burt_width+1)
    membranes =total_burts

    price_df = pd.read_excel('raw_data/price.xlsx')
    price_df.columns = ['name', 'price']
    price_df['name'] = price_df['name'].str.lower()
    price_df['price'] = price_df['price'].replace(39.900000,13.3)
    price_df['price'] = price_df['price'].replace(36833.404400,19000)
    price_df['price'] = price_df['price'].replace(266510.000000,150000)


    price_dict = {each:float(price_df.loc[price_df['name']== each]['price']) for each in list(price_df['name'])}
    price_dict.update({'мембраны': 78878/50*burt_length})

    customs_df = pd.read_excel('raw_data/membrane.xlsx')
    customs_df['Наименование'] = customs_df['Наименование'].str.lower()

    customs_dict = {each:float(customs_df.loc[customs_df['Наименование']==each]['Процент пошлины, %']) for each in list(customs_df['Наименование'])}
    count_dict = {'клей герметик': (round_up((burt_length-3)*2*(1+1/3)))*total_burts+20,
    'блок-пескоуловитель':(2*total_burts),
    'чугунная решётка канала дренажа/вентканала, 150мм':(round_up((burt_length-3)*2*(1+1/3)))*total_burts+20,
    'коммуникационный шкаф с преобразователями':1,
    'кабель слаботочный компьютерный': 20*total_burts,
    'труба гибкая dn110':3*total_burts,
    'клавиатура':1,
    'распределительный шкаф':total_burts,
    'монитор':1,
    'зонд измерения кислорода':total_burts,
    'трубный зажим  dn110':4*total_burts+5,
    'пробки для крышки канала (заглушки отверстий аэрационного канала)': total_burts * 150,
    'датчик давления в комплекте с кабелем':total_burts,
    'держатель настенный для кабеля':total_burts,
    'держатель настенный для 2х зондов':total_burts,
    'защитный козырёк для распределительного шкафа': total_burts,
    'защитный козырёк для вентилятора':total_burts,
    'труба воротник пвх dn110':2* total_burts,
    'труба локтя пвх dn110':2*total_burts,
    'труба редукционная dn160-dn110':2*total_burts,
    'труба прямая пвх':1,
    'труба пвх dn110-110 45°':2*total_burts,
    'труба пвх dn110 2000 мм':2*total_burts+1,
    'труба пвх dn160 1000мм':2*total_burts+1,
    'труба пвх dn110 2000 мм':2*total_burts+1,
    'канал hdpe/pp 150 мм':((burt_length-5)*2*total_burts),
    'конечное звено канала с переходной трубой':2*total_burts*2,
    'заглушка пвх':(2*total_burts),
    'винт-саморез сталь оцинкованная':((round_up((burt_length-3)*2*(1+1/3))))*8*total_burts,
    'кабель слаботочный зонда':total_burts*2+2,
    'зонд температурный':total_burts,
    'вентилятор':total_burts,
    'труба y-образная':total_burts,
    'мембраны':total_burts,
    'намоточная машина':1,
    'ремонтный набор для тента':1
    }
    equipment = list(customs_dict.keys())

    price_no_customs = {each:count_dict[each]*price_dict[each] for each in equipment}
    customs = {each:customs_dict[each] * price_no_customs[each] for each in equipment}
    price_with_customs = {each:customs[each] + price_no_customs[each] for each in equipment}
    price_with_customs_nds = {each:price_with_customs[each] *1.2 for each in equipment}
    price_with_customs_nds_rubl = {each:price_with_customs_nds[each] *euro for each in equipment}
    price_with_customs_nds_rubl.update({'монтаж оборудования':154647*total_burts})
    zagot = ['мембраны','ремонтный набор для тента','кабель слаботочный зонда', 'зонд температурный','вентилятор', 'коммуникационный шкаф с преобразователями', 'кабель слаботочный компьютерный', 'распределительный шкаф', 'зонд измерения кислорода', 'датчик давления в комплекте с кабелем', 'намоточная машина']
    equipment = list(price_with_customs_nds_rubl.keys())
    final_with_nds = {each: (price_with_customs_nds_rubl[each] / 1.2*1.03*1.012*1.2 if each in zagot else price_with_customs_nds_rubl[each] /1.2*1.02*1.2) for each in equipment}
    final_with_nds.update({'монтаж оборудования':154647*total_burts*1.2})
    total_price = sum(final_with_nds.values())

    pres = Presentation('raw_data/base.pptx')
    slides = [slide for slide in pres.slides]
    slide0 = slides[0]
    title_shapes = [shape for shape in slide0.shapes if shape.has_text_frame]
    title = [shape for shape in title_shapes if shape.has_text_frame and shape.text == 'Мощностью {{ tonns }} тонн в год']
    title[0].text = "мощностью %s тонн в год" % ("{:,d}".format(int(tonns)))
    slide2 = pres.slides[2]
    table = [shape for shape in slide2.shapes if shape.has_table]
    table[1].table.cell(1,1).text = "%s" % ("{:,d}".format(int(tonns)))
    class bcolors:
        HEADER = '\033[95m'
        OKBLUE = '\033[94m'
        OKCYAN = '\033[96m'
        OKGREEN = '\033[92m'
        WARNING = '\033[93m'
        FAIL = '\033[91m'
        ENDC = '\033[0m'
        BOLD = '\033[1m'
        UNDERLINE = '\033[4m'
    table[1].table.cell(2,1).text = "%s" % (density)
    table[1].table.cell(3,1).text = "%s недель" % (weeks)
    table[2].table.cell(0,1).text = "%s шт" % (membranes) #буртя 1
    table[2].table.cell(1,1).text = "%s шт" % (total_burts)
    table[4].table.cell(1,1).text = "%s м2" % ("{:,d}".format(int(square)))
    slide7 = slides[5]
    slide7_shapes = [shape for shape in slide7.shapes]
    table = [shape for shape in slide7_shapes if shape.has_table]
    table[0].table.cell(1,1).text  = "%s (шт)\n" % (int(count_dict['канал hdpe/pp 150 мм'] /total_burts))
    table[0].table.cell(1,2).text = "%s (шт)\n" % ("{:,d}".format(int(count_dict['канал hdpe/pp 150 мм'])))
    table[0].table.cell(2,1).text = "%s (шт)\n" % (int(count_dict['чугунная решётка канала дренажа/вентканала, 150мм'] / total_burts))
    table[0].table.cell(2,2).text = "%s (шт)\n" % ("{:,d}".format(int(count_dict['чугунная решётка канала дренажа/вентканала, 150мм'])))
    table[1].table.cell(1,1).text = "%s шт" % (total_burts)
    slide6 = slides[6]
    slide6_shapes = [shape for shape in slide6.shapes]
    table = [shape for shape in slide6_shapes if shape.has_table]
    table[0].table.cell(1,1).text  = "%s" %  total_burts
    table[0].table.cell(2,1).text = "%s" %  total_burts
    table[0].table.cell(3,1).text = "%s" %  total_burts
    table[0].table.cell(4,1).text = "%s" %  total_burts
    table[0].table.cell(5,1).text = "%s" %  total_burts
    table[1].table.cell(2,1).text = 'Установлен на %s\nбуртов' %  total_burts
    #slide11_shapes[8].text_frame.fit_text(bold=False)

    tot_str = "{:,d}".format(int(total_price))
    perc35 = "{:,d}".format(int(total_price*0.35))
    perc60 = "{:,d}".format(int(total_price*0.6))
    perc5 = "{:,d}".format(int(total_price*0.05))
    price_build = 90000*burt_length*total_burts + (total_burts*300000)
    price_build = "{:,d}".format(int(price_build))

    slide11 = slides[11]
    slide11_shapes = [shape for shape in slide11.shapes]

    slide11_shapes[8].text = 'Цена  Договора составляет %s руб. c НДС и включает в себя:\n\nОплата производится по частям в следующем порядке:\n1) Первый платеж в размере 35 процентов от Цены Договора осуществляется после подписания Договора, путем перечисления денежных средств на счет Поставщика и составляет –    %s руб.\n2) Второй платеж в размере 60 процентов от Цены Договора осуществляется после отгрузки Товара в Баден-Бадене, путем перечисления денежных средств на счет Поставщика и составляет –    %s руб.\n3) Третий платеж в размере 5 процентов от Цены Договора осуществляется после ввода оборудования в эксплуатацию, путем перечисления денежных средств на счет Поставщика и составляет –   %s руб.' % (tot_str,perc35,perc60,perc5)

    text_frame = slide11_shapes[8].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Цена  Договора составляет %s руб. c НДС и включает в себя:\n\nОплата производится по частям в следующем порядке:\n1) Первый платеж в размере 35 процентов от Цены Договора осуществляется после подписания Договора, путем перечисления денежных средств на счет Поставщика и составляет –    %s руб.\n2) Второй платеж в размере 60 процентов от Цены Договора осуществляется после отгрузки Товара в Баден-Бадене, путем перечисления денежных средств на счет Поставщика и составляет –    %s руб.\n3) Третий платеж в размере 5 процентов от Цены Договора осуществляется после ввода оборудования в эксплуатацию, путем перечисления денежных средств на счет Поставщика и составляет –   %s руб.' % (tot_str,perc35,perc60,perc5)

    font = run.font
    font.name = 'Geometria Light'
    font.size = Pt(14)
    font.bold = False
    text_frame = slide11_shapes[14].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Строительные работы не входят в поставку, производятся заказчиком самостоятельно и предварительно оцениваются в %s руб.' % (price_build)

    font = run.font
    font.name = 'Geometria Light'
    font.size = Pt(14)
    font.bold = False
    count_dict.update({'монтаж оборудования': 1})
    df = pd.DataFrame.from_dict(final_with_nds,orient='index').reset_index()
    df.columns = ['Наименование', 'Цена с НДС']
    df['Кол-во'] = df['Наименование'].apply(lambda x: count_dict.get(x))
    df['Наименование'] = df['Наименование'].str.capitalize()
    df = df[['Наименование','Кол-во', 'Цена с НДС' ]]

    df.loc['Итого']= df.sum(numeric_only=True, axis=0)
    df['Кол-во'] = df["Кол-во"].astype(int).astype(str)

    df['Цена с НДС'] = df['Цена с НДС'].apply(lambda x: "{:,d}".format(int(x)))
    df.iloc[-1,1] = ''
    df.iloc[-1,0] = 'Итого'
    slide12 = slides[12]
    table = [shape for shape in slide12.shapes if shape.has_table]
    for j in range(0,3):
        table[0].table.cell(0, j).text = df.columns[j]

    for i in range(1,32):
        for j in range(0,3):
          table[0].table.cell(i,j).text = str(df.iloc[i-1, j])
    text = "%s тонн, %s м. длина, %s м. высота стенки, %s недель.pptx"  %("{:,d}".format(int(tonns)),burt_length,burt_wall, weeks)

    print(bcolors.WARNING + text + bcolors.ENDC)

    overview = [{'Буртов':total_burts,
        'Площадь':square,
        'Цена':df.loc[df['Наименование']=='Итого']['Цена с НДС'][0],
        'Рублей за тонну': round_up(sum(final_with_nds.values()) / tonns)

        }
        ]
    overview = pd.DataFrame(overview)
    return pres, overview

if __name__ == '__main__':
    grunteco(10000,50,0.6,5,1, 65)
